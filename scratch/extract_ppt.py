from pptx import Presentation
import os
import sys

# Set encoding for output
sys.stdout.reconfigure(encoding='utf-8')

ppt_path = "ERF_Framework_Presentation (1).pptx"
if os.path.exists(ppt_path):
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
else:
    print("File not found")
